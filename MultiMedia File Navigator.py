import sys
import os
import shutil
import subprocess
import time
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QLabel, QPushButton, 
                             QVBoxLayout, QWidget, QHBoxLayout, QScrollArea, QMessageBox, 
                             QSlider, QStyle)
from PyQt5.QtGui import QPixmap, QImage
from PyQt5.QtCore import Qt, QTimer, QUrl
from PyQt5.QtMultimedia import QMediaPlayer, QMediaContent
from PyQt5.QtMultimediaWidgets import QVideoWidget
import fitz  # PyMuPDF
import traceback
import docx2txt
import openpyxl
import pptx
from odf import opendocument
from odf.text import P

class FileNavigator(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("File Navigator")
        self.setGeometry(100, 100, 1200, 800)

        self.current_file = None
        self.files_to_process = []
        self.destination_folder = None

        self.media_player = QMediaPlayer(None, QMediaPlayer.VideoSurface)
        self.video_widget = QVideoWidget()

        self.init_ui()

    def init_ui(self):
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)

        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.preview_label = QLabel()
        self.preview_label.setAlignment(Qt.AlignCenter)
        self.scroll_area.setWidget(self.preview_label)
        layout.addWidget(self.scroll_area, 1)

        self.file_path_label = QLabel()
        layout.addWidget(self.file_path_label)

        # Media controls
        media_control_layout = QHBoxLayout()
        self.play_pause_button = QPushButton()
        self.play_pause_button.setIcon(self.style().standardIcon(QStyle.SP_MediaPlay))
        self.play_pause_button.clicked.connect(self.toggle_play_pause)
        
        self.timeline_slider = QSlider(Qt.Horizontal)
        self.timeline_slider.sliderMoved.connect(self.set_position)
        
        media_control_layout.addWidget(self.play_pause_button)
        media_control_layout.addWidget(self.timeline_slider)
        layout.addLayout(media_control_layout)

        button_layout = QHBoxLayout()
        self.delete_button = QPushButton("Delete (D)")
        self.move_button = QPushButton("Move (M)")
        self.skip_button = QPushButton("Skip (S)")
        self.open_button = QPushButton("Open (O)")
        button_layout.addWidget(self.delete_button)
        button_layout.addWidget(self.move_button)
        button_layout.addWidget(self.skip_button)
        button_layout.addWidget(self.open_button)
        layout.addLayout(button_layout)

        self.delete_button.clicked.connect(self.delete_file)
        self.move_button.clicked.connect(self.move_file)
        self.skip_button.clicked.connect(self.next_file)
        self.open_button.clicked.connect(self.open_file)

        QTimer.singleShot(0, self.set_shortcuts)

        self.media_player.stateChanged.connect(self.media_state_changed)
        self.media_player.positionChanged.connect(self.position_changed)
        self.media_player.durationChanged.connect(self.duration_changed)

        self.start_navigation()

    def set_shortcuts(self):
        self.delete_button.setShortcut('D')
        self.move_button.setShortcut('M')
        self.skip_button.setShortcut('S')
        self.open_button.setShortcut('O')
        self.play_pause_button.setShortcut('P')

    def start_navigation(self):
        source_path = QFileDialog.getExistingDirectory(self, "Select Source Folder")
        if not source_path:
            self.close()
            return

        self.destination_folder = QFileDialog.getExistingDirectory(self, "Select Destination Folder")
        if not self.destination_folder:
            self.close()
            return

        for root, _, files in os.walk(source_path):
            for file in files:
                self.files_to_process.append(os.path.join(root, file))

        self.next_file()

    def next_file(self):
        self.stop_media_playback()
        if self.files_to_process:
            self.current_file = self.files_to_process.pop(0)
            self.file_path_label.setText(f"Current file: {self.current_file}")
            self.preview_file()
        else:
            QMessageBox.information(self, "Complete", "All files have been processed.")
            self.close()

    def preview_file(self):
        try:
            _, ext = os.path.splitext(self.current_file.lower())
            
            if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                self.preview_image()
            elif ext == '.pdf':
                self.preview_pdf()
            elif ext in ['.doc', '.docx']:
                self.preview_word()
            elif ext == '.odt':
                self.preview_odt()
            elif ext in ['.xls', '.xlsx']:
                self.preview_excel()
            elif ext in ['.ppt', '.pptx']:
                self.preview_powerpoint()
            elif ext in ['.mp4', '.avi', '.mov', '.mp3', '.wav', '.ogg', '.flac']:
                self.preview_media()
            else:
                self.preview_label.setText(f"Preview not available for {ext} files")
            
            self.play_pause_button.setVisible(ext in ['.mp4', '.avi', '.mov', '.mp3', '.wav', '.ogg', '.flac'])
            self.timeline_slider.setVisible(ext in ['.mp4', '.avi', '.mov', '.mp3', '.wav', '.ogg', '.flac'])
        except Exception as e:
            self.handle_preview_error(str(e))

    def preview_image(self):
        try:
            pixmap = QPixmap(self.current_file)
            self.preview_label.setPixmap(pixmap.scaled(1000, 1000, Qt.KeepAspectRatio, Qt.SmoothTransformation))
        except Exception as e:
            raise Exception(f"Error previewing image: {str(e)}")

    def preview_pdf(self):
        try:
            doc = fitz.open(self.current_file)
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format_RGB888)
            pixmap = QPixmap.fromImage(img)
            self.preview_label.setPixmap(pixmap.scaled(1000, 1000, Qt.KeepAspectRatio, Qt.SmoothTransformation))
        except Exception as e:
            raise Exception(f"Error previewing PDF: {str(e)}")

    def preview_word(self):
        try:
            _, ext = os.path.splitext(self.current_file.lower())
            if ext == '.docx':
                content = docx2txt.process(self.current_file)
            else:  # .doc
                content = "Preview not available for .doc files. Use the 'Open' button to view in Notepad."
            
            if not content.strip():
                content = "No text content found in the document."
            self.preview_label.setText(content)
            self.preview_label.setWordWrap(True)
        except Exception as e:
            raise Exception(f"Error previewing Word file: {str(e)}")

    def preview_odt(self):
        try:
            doc = opendocument.load(self.current_file)
            content = []
            
            for element in doc.getElementsByType(P):
                paragraph_text = ''.join(node.data for node in element.childNodes if node.nodeType == 3)  # TEXT_NODE
                if paragraph_text.strip():
                    content.append(paragraph_text)

            if not content:
                content.append("No text content found in the document.")
            
            self.preview_label.setText("\n\n".join(content))
            self.preview_label.setWordWrap(True)
        except Exception as e:
            raise Exception(f"Error previewing ODT file: {str(e)}")

    def preview_excel(self):
        try:
            workbook = openpyxl.load_workbook(self.current_file, data_only=True)
            content = []
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                content.append(f"Sheet: {sheet}")
                for row in worksheet.iter_rows(values_only=True):
                    content.append("\t".join([str(cell) for cell in row if cell is not None]))
            self.preview_label.setText("\n".join(content[:100]))  # Limit to first 100 lines
            self.preview_label.setWordWrap(True)
        except Exception as e:
            raise Exception(f"Error previewing Excel file: {str(e)}")

    def preview_powerpoint(self):
        try:
            prs = pptx.Presentation(self.current_file)
            content = []
            for i, slide in enumerate(prs.slides):
                content.append(f"Slide {i + 1}")
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        content.append(shape.text)
            self.preview_label.setText("\n".join(content))
            self.preview_label.setWordWrap(True)
        except Exception as e:
            raise Exception(f"Error previewing PowerPoint file: {str(e)}")

    def preview_media(self):
        try:
            _, ext = os.path.splitext(self.current_file.lower())
            
            if ext in ['.mp3', '.wav', '.ogg', '.flac']:
                self.scroll_area.setWidget(self.preview_label)
                self.preview_label.setText(f"Audio File ({ext})\nUse Play/Pause button and timeline to control playback")
            else:
                self.scroll_area.setWidget(self.video_widget)
                self.video_widget.show()
            
            self.media_player.setMedia(QMediaContent(QUrl.fromLocalFile(self.current_file)))
        except Exception as e:
            raise Exception(f"Error previewing media file: {str(e)}")

    def toggle_play_pause(self):
        if self.media_player.state() == QMediaPlayer.PlayingState:
            self.media_player.pause()
        else:
            self.media_player.play()

    def media_state_changed(self, state):
        if self.media_player.state() == QMediaPlayer.PlayingState:
            self.play_pause_button.setIcon(self.style().standardIcon(QStyle.SP_MediaPause))
        else:
            self.play_pause_button.setIcon(self.style().standardIcon(QStyle.SP_MediaPlay))

    def position_changed(self, position):
        self.timeline_slider.setValue(position)

    def duration_changed(self, duration):
        self.timeline_slider.setRange(0, duration)

    def set_position(self, position):
        self.media_player.setPosition(position)

    def handle_preview_error(self, error_message):
        self.preview_label.setText(f"Error previewing file:\n{error_message}")
        print(f"Preview error for {self.current_file}: {error_message}")
        print(traceback.format_exc())

    def move_file(self):
        self.stop_media_playback()
        try:
            base_name = os.path.basename(self.current_file)
            name, ext = os.path.splitext(base_name)
            counter = 1
            new_path = os.path.join(self.destination_folder, base_name)

            while os.path.exists(new_path):
                new_name = f"{name}({counter}){ext}"
                new_path = os.path.join(self.destination_folder, new_name)
                counter += 1

            shutil.copy2(self.current_file, new_path)
            
            try:
                os.remove(self.current_file)
            except PermissionError:
                self.force_delete_file()
            
            self.next_file()
            return True
        except PermissionError:
            self.force_delete_file()
            return False
        except Exception as e:
            QMessageBox.warning(self, "Error", f"Failed to move file: {str(e)}")
            return False

    def delete_file(self):
        self.stop_media_playback()
        try:
            os.remove(self.current_file)
            self.next_file()
        except PermissionError:
            self.force_delete_file()
        except Exception as e:
            QMessageBox.warning(self, "Error", f"Failed to delete file: {str(e)}")

    def force_delete_file(self):
        self.stop_media_playback()
        try:
            normalized_path = os.path.normpath(self.current_file)
            escaped_path = normalized_path.replace(" ", "^ ")
            command = f'del /F /Q "{escaped_path}"'
            
            result = subprocess.run(command, 
                                    check=True, capture_output=True, text=True, shell=True)
            
            if result.returncode == 0:
                self.next_file()
            else:
                raise subprocess.CalledProcessError(result.returncode, command, result.stdout, result.stderr)
        
        except subprocess.CalledProcessError as e:
            error_message = f"Failed to force delete file:\nCommand: {e.cmd}\nOutput: {e.stdout}\nError: {e.stderr}"
            QMessageBox.warning(self, "Error", error_message)
        except Exception as e:
            QMessageBox.warning(self, "Error", f"An unexpected error occurred: {str(e)}")

    def stop_media_playback(self):
        if self.media_player.state() == QMediaPlayer.PlayingState:
            self.media_player.stop()

    def open_file(self):
        try:
            _, ext = os.path.splitext(self.current_file.lower())
            if ext == '.doc':
                subprocess.Popen(['notepad.exe', self.current_file])
            else:
                os.startfile(self.current_file)
        except Exception as e:
            QMessageBox.warning(self, "Error", f"Failed to open file: {str(e)}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = FileNavigator()
    window.show()
    sys.exit(app.exec_())